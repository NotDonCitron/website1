#!/usr/bin/env python3
"""
Presentation Builder Module
Automatische PowerPoint-Präsentationserstellung für Trading-Bot Marketing

Features:
- Dynamische Slide-Generierung aus Trade-Daten
- Professionelle Templates mit Animationen
- Charts und Diagramme
- Multi-Format Export (PPTX, PDF, HTML)
- Marketing-optimierte Layouts
"""

import os
import json
import logging
from typing import Dict, List, Optional, Tuple, Any
from datetime import datetime, timedelta
from dataclasses import dataclass
from pathlib import Path

# PowerPoint Libraries
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# Zusätzliche Libraries für Charts und Export
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.backends.backend_pdf import PdfPages
import seaborn as sns
import pandas as pd
import numpy as np
from PIL import Image, ImageDraw, ImageFont
import base64
from io import BytesIO

@dataclass
class SlideContent:
    """Datenklasse für Slide-Inhalte"""
    title: str
    subtitle: str = ""
    content: List[str] = None
    image_path: str = ""
    chart_data: Dict = None
    layout_type: str = "content"  # title, content, image, chart, comparison
    animation: str = "none"
    background_color: str = "#FFFFFF"

class PresentationBuilder:
    """Hauptklasse für Präsentationserstellung"""
    
    def __init__(self, template_manager=None):
        self.logger = logging.getLogger(__name__)
        self.template_manager = template_manager
        
        # Statistiken für Charts
        self.performance_stats = {}
        
        # Farb-Paletten für verschiedene Themes
        self.color_schemes = {
            'crypto_professional': {
                'primary': RGBColor(26, 35, 126),      # Dunkelblau
                'secondary': RGBColor(63, 81, 181),     # Blau
                'accent': RGBColor(255, 193, 7),        # Gold
                'success': RGBColor(76, 175, 80),       # Grün
                'danger': RGBColor(244, 67, 54),        # Rot
                'text': RGBColor(33, 33, 33),           # Dunkelgrau
                'background': RGBColor(248, 249, 250)   # Hellgrau
            },
            'modern_dark': {
                'primary': RGBColor(18, 18, 18),        # Schwarz
                'secondary': RGBColor(45, 45, 45),      # Dunkelgrau
                'accent': RGBColor(0, 230, 118),        # Neongrün
                'success': RGBColor(0, 200, 83),        # Grün
                'danger': RGBColor(255, 82, 82),        # Rot
                'text': RGBColor(255, 255, 255),        # Weiß
                'background': RGBColor(30, 30, 30)      # Dunkelgrau
            },
            'corporate_blue': {
                'primary': RGBColor(0, 77, 153),        # Corporate Blau
                'secondary': RGBColor(0, 102, 204),     # Hellblau
                'accent': RGBColor(255, 140, 0),        # Orange
                'success': RGBColor(34, 139, 34),       # Grün
                'danger': RGBColor(220, 20, 60),        # Rot
                'text': RGBColor(51, 51, 51),           # Grau
                'background': RGBColor(255, 255, 255)   # Weiß
            }
        }
        
    def create_presentation(self, trades_data: List[Dict], config: Dict, 
                          template: str = "crypto_professional", 
                          output_path: str = "presentation", 
                          format: str = "pptx") -> bool:
        """Hauptfunktion für Präsentationserstellung"""
        try:
            self.logger.info(f"Erstelle Präsentation mit {len(trades_data)} Trades")
            
            # PowerPoint-Präsentation erstellen
            prs = Presentation()
            
            # Template-spezifische Einstellungen
            colors = self.color_schemes.get(template, self.color_schemes['crypto_professional'])
            
            # Slides generieren
            self._create_title_slide(prs, config, colors)
            self._create_overview_slide(prs, config, trades_data, colors)
            self._create_statistics_slide(prs, config, trades_data, colors)
            
            # Trade-spezifische Slides
            for i, trade in enumerate(trades_data[:20]):  # Maximal 20 Trades
                self._create_trade_slide(prs, trade, colors, i + 1)
                
            # Performance-Charts
            if config.get('presentation', {}).get('include_charts', True):
                self._create_performance_chart_slide(prs, trades_data, colors)
                self._create_roi_distribution_slide(prs, trades_data, colors)
                
            # Call-to-Action Slides
            if config.get('presentation', {}).get('include_cta', True):
                self._create_cta_slide(prs, config, colors)
                self._create_contact_slide(prs, config, colors)
                
            # Backup erstellen
            if config.get('processing', {}).get('backup_enabled', True):
                self._create_backup(trades_data, config, output_path)
                
            # Export in gewünschtem Format
            success = self._export_presentation(prs, output_path, format, trades_data)
            
            if success:
                self.logger.info(f"Präsentation erfolgreich erstellt: {output_path}.{format}")
                return True
            else:
                return False
                
        except Exception as e:
            self.logger.error(f"Fehler beim Erstellen der Präsentation: {e}")
            return False
            
    def _create_title_slide(self, prs: Presentation, config: Dict, colors: Dict):
        """Titel-Slide erstellen"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Titel
        title = slide.shapes.title
        title.text = config.get('presentation', {}).get('title', '🤖 Profitabler Trading Bot')
        
        # Titel formatieren
        title_para = title.text_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = colors['primary']
        title_para.alignment = PP_ALIGN.CENTER
        
        # Untertitel
        subtitle = slide.placeholders[1]
        bot_data = config.get('bot_data', {})
        weekly_stats = bot_data.get('weekly_stats', {})
        
        subtitle_text = f"""Automatischer Telegram-Signal Bot für Bybit
        
📊 Wochenstatistik ({weekly_stats.get('period', 'N/A')}):
• {weekly_stats.get('trades', 0)} Trades ausgeführt
• {weekly_stats.get('winrate', 0)}% Erfolgsrate
• +{weekly_stats.get('roi', 0)}% ROI
        
🚀 Referral Code: {bot_data.get('referral_code', 'N/A')}"""
        
        subtitle.text = subtitle_text
        
        # Untertitel formatieren
        for paragraph in subtitle.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = colors['text']
            paragraph.alignment = PP_ALIGN.CENTER
            
    def _create_overview_slide(self, prs: Presentation, config: Dict, 
                             trades_data: List[Dict], colors: Dict):
        """Überblick-Slide erstellen"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Titel
        title = slide.shapes.title
        title.text = "📈 Trading Bot Übersicht"
        self._format_title(title, colors)
        
        # Content
        content = slide.placeholders[1]
        
        # Statistiken berechnen
        total_trades = len(trades_data)
        winning_trades = len([t for t in trades_data if t.get('status') == 'win'])
        losing_trades = len([t for t in trades_data if t.get('status') == 'loss'])
        pending_trades = len([t for t in trades_data if t.get('status') == 'pending'])
        
        roi_values = [t.get('roi', 0) for t in trades_data if t.get('roi') is not None]
        avg_roi = sum(roi_values) / len(roi_values) if roi_values else 0
        max_roi = max(roi_values) if roi_values else 0
        
        content_text = f"""✅ WAS MACHT DIESER BOT?
• Automatisches Kopieren von Telegram-Signalen
• Direkte Ausführung auf Bybit Exchange
• 24/7 Trading ohne manuellen Eingriff
• Intelligente Risk-Management Features

📊 AKTUELLE PERFORMANCE:
• Gesamte Trades: {total_trades}
• Gewinn-Trades: {winning_trades} ({(winning_trades/total_trades*100):.1f}%)
• Verlust-Trades: {losing_trades} ({(losing_trades/total_trades*100):.1f}%)
• Durchschnittlicher ROI: +{avg_roi:.2f}%
• Bester Trade: +{max_roi:.2f}%

🎯 WARUM FUNKTIONIERT ES?
• Bewährte Telegram-Signal Quellen
• Schnelle Ausführung ohne Verzögerung
• Automatisches Stop-Loss Management
• Kontinuierliche Performance-Optimierung"""
        
        content.text = content_text
        self._format_content(content, colors)
        
    def _create_statistics_slide(self, prs: Presentation, config: Dict, 
                               trades_data: List[Dict], colors: Dict):
        """Statistik-Slide mit Highlights erstellen"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Titel
        title = slide.shapes.title
        title.text = "🏆 Top Performer & Statistiken"
        self._format_title(title, colors)
        
        # Top Performer ermitteln
        sorted_trades = sorted(trades_data, key=lambda x: x.get('roi', 0), reverse=True)
        top_trades = sorted_trades[:5]
        
        # Bot-Daten aus Config
        bot_data = config.get('bot_data', {})
        top_performers = bot_data.get('top_performers', [])
        example_trades = bot_data.get('example_trades', [])
        
        content = slide.placeholders[1]
        
        content_text = f"""🥇 TOP PERFORMING COINS:"""
        
        # Top Performer aus Config anzeigen
        for i, performer in enumerate(top_performers[:3], 1):
            content_text += f"\n{i}. {performer.get('coin', 'N/A')}: +{performer.get('roi', 0)}% ROI"
            
        content_text += f"\n\n📈 BEISPIEL-TRADES:"
        
        # Beispiel-Trades aus aktuellen Daten
        for trade in top_trades[:5]:
            coin = trade.get('coin', 'N/A')
            roi = trade.get('roi', 0)
            status_emoji = "✅" if trade.get('status') == 'win' else "❌" if trade.get('status') == 'loss' else "⏳"
            content_text += f"\n{status_emoji} {coin}: "
            if roi > 0:
                content_text += f"+{roi:.2f}%"
            else:
                content_text += f"{roi:.2f}%"
                
        # Zusätzliche Statistiken
        roi_values = [t.get('roi', 0) for t in trades_data if t.get('roi') is not None and t.get('roi') > 0]
        if roi_values:
            content_text += f"\n\n📊 WEITERE STATISTIKEN:"
            content_text += f"\n• Durchschnittlicher Gewinn: +{sum(roi_values)/len(roi_values):.2f}%"
            content_text += f"\n• Median ROI: +{sorted(roi_values)[len(roi_values)//2]:.2f}%"
            content_text += f"\n• Trades über +100%: {len([r for r in roi_values if r > 100])}"
            content_text += f"\n• Trades über +500%: {len([r for r in roi_values if r > 500])}"
            
        content.text = content_text
        self._format_content(content, colors)
        
    def _create_trade_slide(self, prs: Presentation, trade: Dict, colors: Dict, slide_num: int):
        """Einzelnen Trade-Slide erstellen"""
        slide_layout = prs.slide_layouts[5]  # Blank layout für mehr Kontrolle
        slide = prs.slides.add_slide(slide_layout)
        
        # Titel hinzufügen
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        
        coin = trade.get('coin', 'Unknown')
        roi = trade.get('roi', 0)
        status = trade.get('status', 'unknown')
        
        # Status-Emoji
        status_emoji = "✅" if status == 'win' else "❌" if status == 'loss' else "⏳"
        
        title_text = f"{status_emoji} {coin} Trade #{slide_num}"
        if roi is not None:
            if roi > 0:
                title_text += f" (+{roi:.2f}%)"
            else:
                title_text += f" ({roi:.2f}%)"
                
        title_frame.text = title_text
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = colors['primary']
        title_para.alignment = PP_ALIGN.CENTER
        
        # Haupt-Content Box
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
        content_frame = content_box.text_frame
        
        # Trade-Details zusammenstellen
        details = []
        
        if trade.get('entry_price'):
            details.append(f"📍 Entry Price: ${trade.get('entry_price'):.6f}")
            
        if trade.get('exit_price'):
            details.append(f"🎯 Exit Price: ${trade.get('exit_price'):.6f}")
            
        if roi is not None:
            roi_color = "🟢" if roi > 0 else "🔴" if roi < 0 else "🟡"
            details.append(f"{roi_color} ROI: {roi:+.2f}%")
            
        if trade.get('timestamp'):
            try:
                timestamp = datetime.fromisoformat(trade.get('timestamp'))
                details.append(f"⏰ Timestamp: {timestamp.strftime('%d.%m.%Y %H:%M')}")
            except:
                pass
                
        # Signal/Result Images
        if trade.get('signal_image'):
            details.append(f"📊 Signal Screenshot: Verfügbar")
            
        if trade.get('result_image'):
            details.append(f"💰 Ergebnis Screenshot: Verfügbar")
            
        match_confidence = trade.get('match_confidence', 0)
        if match_confidence > 0:
            details.append(f"🎲 Match Confidence: {match_confidence:.1%}")
            
        content_text = "\n\n".join(details)
        
        # Zusätzliche Bewertung basierend auf ROI
        if roi is not None:
            content_text += "\n\n" + "="*40 + "\n"
            
            if roi > 500:
                content_text += "🚀 EXCEPTIONAL TRADE!\nDieser Trade zeigt das enorme Potential des Bots."
            elif roi > 200:
                content_text += "⭐ EXCELLENT PERFORMANCE!\nSolche Trades machen den Bot profitabel."
            elif roi > 50:
                content_text += "✨ SOLID PROFIT!\nKonsistente Gewinne wie dieser sind das Ziel."
            elif roi > 0:
                content_text += "✅ PROFITABLE TRADE\nJeder Gewinn zählt für die Gesamtperformance."
            elif roi < -50:
                content_text += "⚠️ LEARNING EXPERIENCE\nSolche Trades sind selten und Teil des Lernprozesses."
            else:
                content_text += "📊 ANALYZED TRADE\nJeder Trade wird dokumentiert für Transparenz."
                
        content_frame.text = content_text
        
        # Content formatieren
        for para in content_frame.paragraphs:
            para.font.size = Pt(16)
            para.font.color.rgb = colors['text']
            para.alignment = PP_ALIGN.LEFT
            
        # ROI-abhängige Hintergrundfarbe (subtil)
        if roi is not None:
            if roi > 100:
                # Grüner Hintergrund für sehr gute Trades
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(240, 255, 240)
            elif roi < -20:
                # Leicht rötlicher Hintergrund für Verluste
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(255, 248, 248)
                
    def _create_performance_chart_slide(self, prs: Presentation, trades_data: List[Dict], colors: Dict):
        """Performance-Chart Slide erstellen"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Titel
        title = slide.shapes.title
        title.text = "📊 Performance Analyse"
        self._format_title(title, colors)
        
        try:
            # Chart-Daten vorbereiten
            roi_values = [t.get('roi', 0) for t in trades_data if t.get('roi') is not None]
            
            if roi_values:
                # Matplotlib Chart erstellen
                fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 6))
                
                # ROI Histogram
                ax1.hist(roi_values, bins=20, alpha=0.7, color='steelblue', edgecolor='black')
                ax1.set_title('ROI Verteilung')
                ax1.set_xlabel('ROI (%)')
                ax1.set_ylabel('Anzahl Trades')
                ax1.grid(True, alpha=0.3)
                
                # ROI über Zeit (falls Timestamps verfügbar)
                timestamps = []
                for trade in trades_data:
                    if trade.get('timestamp') and trade.get('roi') is not None:
                        try:
                            ts = datetime.fromisoformat(trade.get('timestamp'))
                            timestamps.append((ts, trade.get('roi')))
                        except:
                            continue
                            
                if timestamps:
                    timestamps.sort(key=lambda x: x[0])
                    dates, rois = zip(*timestamps)
                    
                    ax2.plot(dates, rois, marker='o', linewidth=2, markersize=4)
                    ax2.set_title('ROI über Zeit')
                    ax2.set_xlabel('Zeit')
                    ax2.set_ylabel('ROI (%)')
                    ax2.grid(True, alpha=0.3)
                    plt.setp(ax2.xaxis.get_majorticklabels(), rotation=45)
                else:
                    # Fallback: Top Trades Bar Chart
                    sorted_trades = sorted(trades_data, key=lambda x: x.get('roi', 0), reverse=True)[:10]
                    coins = [t.get('coin', f'Trade {i}') for i, t in enumerate(sorted_trades)]
                    rois = [t.get('roi', 0) for t in sorted_trades]
                    
                    bars = ax2.bar(coins, rois, color=['green' if r > 0 else 'red' for r in rois])
                    ax2.set_title('Top 10 Trades')
                    ax2.set_xlabel('Coin')
                    ax2.set_ylabel('ROI (%)')
                    plt.setp(ax2.xaxis.get_majorticklabels(), rotation=45)
                    
                plt.tight_layout()
                
                # Chart als Bild speichern und einfügen
                chart_path = self._save_chart_as_image(fig, "performance_chart")
                if chart_path and os.path.exists(chart_path):
                    # Chart in Slide einfügen
                    slide.shapes.add_picture(chart_path, Inches(1), Inches(2), 
                                           width=Inches(8), height=Inches(4))
                    
                plt.close(fig)
                
        except Exception as e:
            self.logger.error(f"Fehler beim Erstellen des Performance-Charts: {e}")
            
            # Fallback: Text-basierte Statistiken
            content = slide.placeholders[1]
            content_text = self._generate_text_statistics(trades_data)
            content.text = content_text
            self._format_content(content, colors)
            
    def _create_roi_distribution_slide(self, prs: Presentation, trades_data: List[Dict], colors: Dict):
        """ROI-Verteilung Slide erstellen"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Titel
        title = slide.shapes.title
        title.text = "💰 ROI Verteilung & Kategorien"
        self._format_title(title, colors)
        
        # ROI-Kategorien analysieren
        roi_values = [t.get('roi', 0) for t in trades_data if t.get('roi') is not None]
        
        categories = {
            'Moonshots (>1000%)': len([r for r in roi_values if r > 1000]),
            'Große Gewinne (500-1000%)': len([r for r in roi_values if 500 <= r <= 1000]),
            'Sehr gut (200-500%)': len([r for r in roi_values if 200 <= r < 500]),
            'Gut (100-200%)': len([r for r in roi_values if 100 <= r < 200]),
            'Profitabel (0-100%)': len([r for r in roi_values if 0 < r < 100]),
            'Break-even (0%)': len([r for r in roi_values if r == 0]),
            'Verluste (<0%)': len([r for r in roi_values if r < 0])
        }
        
        content = slide.placeholders[1]
        
        content_text = "📈 ROI KATEGORIEN:\n\n"
        
        total_trades = len(roi_values)
        for category, count in categories.items():
            if count > 0:
                percentage = (count / total_trades) * 100
                content_text += f"• {category}: {count} Trades ({percentage:.1f}%)\n"
                
        # Zusätzliche Statistiken
        if roi_values:
            content_text += f"\n📊 STATISTISCHE KENNZAHLEN:\n"
            content_text += f"• Durchschnitt: {np.mean(roi_values):.2f}%\n"
            content_text += f"• Median: {np.median(roi_values):.2f}%\n"
            content_text += f"• Standardabweichung: {np.std(roi_values):.2f}%\n"
            content_text += f"• Minimum: {min(roi_values):.2f}%\n"
            content_text += f"• Maximum: {max(roi_values):.2f}%\n"
            
            # Erfolgsrate
            winning_trades = len([r for r in roi_values if r > 0])
            win_rate = (winning_trades / total_trades) * 100
            content_text += f"\n🎯 ERFOLGSRATE: {win_rate:.1f}%\n"
            
            # Risk-Reward Verhältnis
            avg_win = np.mean([r for r in roi_values if r > 0]) if winning_trades > 0 else 0
            losing_trades = [r for r in roi_values if r < 0]
            avg_loss = abs(np.mean(losing_trades)) if losing_trades else 0
            
            if avg_loss > 0:
                risk_reward = avg_win / avg_loss
                content_text += f"⚖️ Risk-Reward Ratio: {risk_reward:.2f}:1"
                
        content.text = content_text
        self._format_content(content, colors)
        
    def _create_cta_slide(self, prs: Presentation, config: Dict, colors: Dict):
        """Call-to-Action Slide erstellen"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Titel
        title = slide.shapes.title
        title.text = "🚀 Bereit für automatische Gewinne?"
        self._format_title(title, colors)
        
        content = slide.placeholders[1]
        referral_code = config.get('bot_data', {}).get('referral_code', 'N/A')
        
        content_text = f"""💎 WARUM SOLLTEN SIE JETZT STARTEN?

✅ Bewiesene Performance mit echten Ergebnissen
✅ 24/7 automatisches Trading - kein manueller Aufwand
✅ Professionelle Signal-Quellen mit hoher Trefferquote
✅ Transparente Dokumentation aller Trades
✅ Kontinuierliche Optimierung der Strategien

🎯 LIMITIERTES ANGEBOT:
Nur die ersten 100 Nutzer erhalten Zugang zu diesen
exklusiven Signal-Quellen und der bewährten Konfiguration.

💰 SO STARTEN SIE:
1. Bybit Account mit Referral Code erstellen: {referral_code}
2. Bot-Setup und Konfiguration erhalten
3. Automatisches Trading startet sofort
4. Gewinne direkt auf Ihr Bybit Konto

⚡ BONUS für Early Adopters:
• Kostenlose Ersteinrichtung
• 30 Tage Premium-Support
• Exklusive Telegram-Gruppe
• Monatliche Performance-Reports

🔥 JETZT HANDELN - PLÄTZE SIND BEGRENZT!"""
        
        content.text = content_text
        self._format_content(content, colors, size=Pt(16))
        
    def _create_contact_slide(self, prs: Presentation, config: Dict, colors: Dict):
        """Kontakt-Slide erstellen"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Titel
        title = slide.shapes.title
        title.text = "📞 Kontakt & Nächste Schritte"
        self._format_title(title, colors)
        
        content = slide.placeholders[1]
        referral_code = config.get('bot_data', {}).get('referral_code', 'N/A')
        
        content_text = f"""📧 KONTAKTIEREN SIE UNS FÜR DEN SOFORTIGEN START:

🔗 Bybit Referral Code: {referral_code}
   (Verwenden Sie diesen Code für optimale Einstellungen)

📱 Telegram: @YourTradingBot
   (Schnellste Antwortzeit - meist innerhalb von 30 Minuten)

💻 E-Mail: info@yourtradingbot.com
   (Für detaillierte Fragen und technischen Support)

🌐 Website: www.yourtradingbot.com
   (Weitere Informationen und Tutorials)

⚡ SOFORT-SETUP VERFÜGBAR:
Nach Ihrer Kontaktaufnahme erhalten Sie innerhalb von 24h:
• Vollständige Bot-Konfiguration
• Schritt-für-Schritt Setup-Anleitung
• Zugang zur privaten Telegram-Gruppe
• Erste Signale und Trades

🎁 AKTUELLE BONI:
• Erste Woche kostenlos testen
• Geld-zurück-Garantie bei Unzufriedenheit
• Persönliche Einrichtungsberatung
• Lebenslanger Support-Zugang

💼 VERTRAUEN SIE AUF BEWIESENE ERGEBNISSE!
Alle gezeigten Trades sind authentisch und nachweisbar.
Ihre Investition ist durch unsere Erfolgsbilanz gesichert."""
        
        content.text = content_text
        self._format_content(content, colors, size=Pt(14))
        
    def _format_title(self, title_shape, colors: Dict):
        """Titel-Formatierung standardisieren"""
        title_frame = title_shape.text_frame
        for paragraph in title_frame.paragraphs:
            paragraph.font.size = Pt(36)
            paragraph.font.bold = True
            paragraph.font.color.rgb = colors['primary']
            paragraph.alignment = PP_ALIGN.CENTER
            
    def _format_content(self, content_shape, colors: Dict, size: Pt = Pt(18)):
        """Content-Formatierung standardisieren"""
        content_frame = content_shape.text_frame
        for paragraph in content_frame.paragraphs:
            paragraph.font.size = size
            paragraph.font.color.rgb = colors['text']
            paragraph.alignment = PP_ALIGN.LEFT
            
    def _save_chart_as_image(self, fig, filename: str) -> str:
        """Matplotlib Chart als Bild speichern"""
        try:
            chart_path = f"{filename}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png"
            fig.savefig(chart_path, dpi=300, bbox_inches='tight', 
                       facecolor='white', edgecolor='none')
            return chart_path
        except Exception as e:
            self.logger.error(f"Fehler beim Speichern des Charts: {e}")
            return ""
            
    def _generate_text_statistics(self, trades_data: List[Dict]) -> str:
        """Text-basierte Statistiken als Fallback"""
        roi_values = [t.get('roi', 0) for t in trades_data if t.get('roi') is not None]
        
        if not roi_values:
            return "Keine ROI-Daten verfügbar für Statistiken."
            
        stats_text = f"""📊 PERFORMANCE STATISTIKEN:

Gesamte Trades: {len(roi_values)}
Durchschnittlicher ROI: {np.mean(roi_values):.2f}%
Median ROI: {np.median(roi_values):.2f}%
Bester Trade: {max(roi_values):.2f}%
Schlechtester Trade: {min(roi_values):.2f}%

Gewinn-Trades: {len([r for r in roi_values if r > 0])}
Verlust-Trades: {len([r for r in roi_values if r < 0])}
Erfolgsrate: {(len([r for r in roi_values if r > 0]) / len(roi_values) * 100):.1f}%

Trades über +100%: {len([r for r in roi_values if r > 100])}
Trades über +500%: {len([r for r in roi_values if r > 500])}
Trades über +1000%: {len([r for r in roi_values if r > 1000])}"""
        
        return stats_text
        
    def _export_presentation(self, prs: Presentation, output_path: str, 
                           format: str, trades_data: List[Dict]) -> bool:
        """Präsentation in verschiedenen Formaten exportieren"""
        try:
            if format.lower() == 'pptx':
                prs.save(f"{output_path}.pptx")
                return True
                
            elif format.lower() == 'pdf':
                return self._export_to_pdf(prs, f"{output_path}.pdf", trades_data)
                
            elif format.lower() == 'html':
                return self._export_to_html(trades_data, f"{output_path}.html")
                
            else:
                self.logger.error(f"Unbekanntes Export-Format: {format}")
                return False
                
        except Exception as e:
            self.logger.error(f"Fehler beim Export: {e}")
            return False
            
    def _export_to_pdf(self, prs: Presentation, output_path: str, trades_data: List[Dict]) -> bool:
        """PDF-Export (vereinfacht)"""
        try:
            # Für PDF-Export würde man normalerweise die Slides als Bilder rendern
            # und dann in ein PDF zusammenfügen. Hier eine vereinfachte Version:
            
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            
            doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            # Titel-Seite
            title = Paragraph("Trading Bot Performance Report", styles['Title'])
            story.append(title)
            story.append(Spacer(1, 12))
            
            # Statistiken hinzufügen
            stats_text = self._generate_text_statistics(trades_data)
            stats_para = Paragraph(stats_text.replace('\n', '<br/>'), styles['Normal'])
            story.append(stats_para)
            
            # Trade-Details
            for i, trade in enumerate(trades_data[:10]):  # Erste 10 Trades
                story.append(Spacer(1, 12))
                trade_title = f"Trade {i+1}: {trade.get('coin', 'N/A')}"
                trade_para = Paragraph(trade_title, styles['Heading2'])
                story.append(trade_para)
                
                trade_details = f"ROI: {trade.get('roi', 0):.2f}% | Status: {trade.get('status', 'N/A')}"
                details_para = Paragraph(trade_details, styles['Normal'])
                story.append(details_para)
                
            doc.build(story)
            return True
            
        except ImportError:
            self.logger.error("ReportLab nicht installiert - PDF-Export nicht verfügbar")
            return False
        except Exception as e:
            self.logger.error(f"Fehler beim PDF-Export: {e}")
            return False
            
    def _export_to_html(self, trades_data: List[Dict], output_path: str) -> bool:
        """HTML-Export für Web-Präsentation"""
        try:
            html_content = self._generate_html_presentation(trades_data)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
                
            return True
            
        except Exception as e:
            self.logger.error(f"Fehler beim HTML-Export: {e}")
            return False
            
    def _generate_html_presentation(self, trades_data: List[Dict]) -> str:
        """HTML-Präsentation generieren"""
        # ROI-Statistiken
        roi_values = [t.get('roi', 0) for t in trades_data if t.get('roi') is not None]
        avg_roi = np.mean(roi_values) if roi_values else 0
        max_roi = max(roi_values) if roi_values else 0
        win_rate = (len([r for r in roi_values if r > 0]) / len(roi_values) * 100) if roi_values else 0
        
        html = f"""
<!DOCTYPE html>
<html lang="de">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Trading Bot Performance Report</title>
    <style>
        body {{
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            margin: 0;
            padding: 20px;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: #333;
        }}
        .container {{
            max-width: 1200px;
            margin: 0 auto;
            background: white;
            border-radius: 15px;
            padding: 30px;
            box-shadow: 0 10px 30px rgba(0,0,0,0.3);
        }}
        .header {{
            text-align: center;
            padding: 20px 0;
            border-bottom: 3px solid #667eea;
            margin-bottom: 30px;
        }}
        .header h1 {{
            color: #667eea;
            font-size: 2.5em;
            margin: 0;
        }}
        .stats-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(250px, 1fr));
            gap: 20px;
            margin: 30px 0;
        }}
        .stat-card {{
            background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%);
            padding: 20px;
            border-radius: 10px;
            text-align: center;
            color: white;
        }}
        .stat-value {{
            font-size: 2em;
            font-weight: bold;
            margin: 10px 0;
        }}
        .trades-section {{
            margin: 30px 0;
        }}
        .trade-card {{
            background: #f8f9fa;
            border-left: 5px solid #28a745;
            padding: 15px;
            margin: 10px 0;
            border-radius: 5px;
        }}
        .trade-card.loss {{
            border-left-color: #dc3545;
        }}
        .trade-header {{
            font-weight: bold;
            font-size: 1.2em;
            margin-bottom: 5px;
        }}
        .roi-positive {{
            color: #28a745;
            font-weight: bold;
        }}
        .roi-negative {{
            color: #dc3545;
            font-weight: bold;
        }}
        .cta-section {{
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white;
            padding: 30px;
            border-radius: 10px;
            text-align: center;
            margin: 30px 0;
        }}
        .cta-button {{
            background: #ffc107;
            color: #333;
            padding: 15px 30px;
            border: none;
            border-radius: 25px;
            font-size: 1.2em;
            font-weight: bold;
            cursor: pointer;
            margin: 10px;
        }}
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1>🤖 Trading Bot Performance Report</h1>
            <p>Automatischer Telegram-Signal Bot für Bybit</p>
        </div>
        
        <div class="stats-grid">
            <div class="stat-card">
                <div>Gesamte Trades</div>
                <div class="stat-value">{len(trades_data)}</div>
            </div>
            <div class="stat-card">
                <div>Durchschnittlicher ROI</div>
                <div class="stat-value">{avg_roi:.1f}%</div>
            </div>
            <div class="stat-card">
                <div>Bester Trade</div>
                <div class="stat-value">{max_roi:.1f}%</div>
            </div>
            <div class="stat-card">
                <div>Erfolgsrate</div>
                <div class="stat-value">{win_rate:.1f}%</div>
            </div>
        </div>
        
        <div class="trades-section">
            <h2>🏆 Top Performing Trades</h2>
"""
        
        # Top 10 Trades hinzufügen
        sorted_trades = sorted(trades_data, key=lambda x: x.get('roi', 0), reverse=True)[:10]
        
        for i, trade in enumerate(sorted_trades, 1):
            coin = trade.get('coin', 'N/A')
            roi = trade.get('roi', 0)
            status = trade.get('status', 'unknown')
            
            roi_class = 'roi-positive' if roi > 0 else 'roi-negative'
            card_class = 'trade-card' if roi > 0 else 'trade-card loss'
            status_emoji = "✅" if status == 'win' else "❌" if status == 'loss' else "⏳"
            
            html += f"""
            <div class="{card_class}">
                <div class="trade-header">{status_emoji} #{i} {coin}</div>
                <div>ROI: <span class="{roi_class}">{roi:+.2f}%</span></div>
            </div>
"""
        
        html += """
        </div>
        
        <div class="cta-section">
            <h2>🚀 Bereit für automatische Gewinne?</h2>
            <p>Starten Sie noch heute mit dem bewährten Trading Bot!</p>
            <button class="cta-button">Jetzt starten</button>
            <button class="cta-button">Mehr Informationen</button>
        </div>
    </div>
</body>
</html>"""
        
        return html
        
    def _create_backup(self, trades_data: List[Dict], config: Dict, output_path: str):
        """Backup der Daten erstellen"""
        try:
            backup_data = {
                'timestamp': datetime.now().isoformat(),
                'config': config,
                'trades_data': trades_data,
                'metadata': {
                    'total_trades': len(trades_data),
                    'export_path': output_path,
                    'version': '1.0.0'
                }
            }
            
            backup_path = f"{output_path}_backup_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
            
            with open(backup_path, 'w', encoding='utf-8') as f:
                json.dump(backup_data, f, indent=2, ensure_ascii=False)
                
            self.logger.info(f"Backup erstellt: {backup_path}")
            
        except Exception as e:
            self.logger.error(f"Fehler beim Erstellen des Backups: {e}")

if __name__ == "__main__":
    # Test-Code
    test_trades = [
        {'coin': 'BTC', 'roi': 150.5, 'status': 'win'},
        {'coin': 'ETH', 'roi': -25.2, 'status': 'loss'},
        {'coin': 'EPIC', 'roi': 5261.0, 'status': 'win'},
    ]
    
    test_config = {
        'bot_data': {
            'referral_code': '9EEDV9',
            'weekly_stats': {'trades': 59, 'winrate': 84.7, 'roi': 31.690}
        },
        'presentation': {
            'title': 'Test Präsentation',
            'include_charts': True,
            'include_cta': True
        }
    }
    
    builder = PresentationBuilder()
    success = builder.create_presentation(test_trades, test_config, output_path="test_presentation")
    
    if success:
        print("✅ Test-Präsentation erfolgreich erstellt!")
    else:
        print("❌ Fehler beim Erstellen der Test-Präsentation")