
"""
🚀 AUTOMATISCHER TRADING BOT PRÄSENTATIONS-GENERATOR
Erstellt automatisch eine PowerPoint-Präsentation mit echten Trading-Beispielen

ANLEITUNG:
1. Installiere benötigte Pakete: pip install python-pptx pillow requests
2. Lege alle deine Bilder in einen Ordner namens "trading_images"
3. Führe dieses Skript aus
4. Fertige Präsentation wird als "Trading_Bot_Final.pptx" gespeichert
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
import re
from PIL import Image
import requests
from io import BytesIO

def create_trading_presentation():
    # Neue Präsentation erstellen
    prs = Presentation()

    # === SLIDE 1: TITEL ===
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "🚀 AUTOMATISCHER TRADING BOT"
    subtitle.text = "Telegram Signal Kopierer\nVollautomatische Gewinnmaximierung\n\nErgebnisse: 17.07 - 22.07.2025"

    # === SLIDE 2: FUNKTIONSWEISE ===
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "🤖 WIE FUNKTIONIERT DER BOT?"
    content.text = """• Überwacht automatisch Telegram-Signalgruppen 24/7
• Kopiert Handelssignale in Echtzeit
• Führt Trades automatisch auf Bybit aus
• Verwendet optimierte Hebel (12.5X - 50X)
• Setzt automatische Take-Profit Ziele
• Kein manueller Eingriff erforderlich

✅ Vollautomatisch
✅ 24/7 aktiv  
✅ Sofortige Ausführung
✅ Risikomanagement integriert"""

    # === SLIDE 3: STATISTIKEN ===
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "📊 WOCHENSTATISTIKEN (17.07 - 22.07)"
    content.text = """🎯 GESAMT TRADES: 59

🟢 GEWINN-TRADES: 50 (84.7% Winrate!)
🟡 BREAKEVEN: 4 
🔴 VERLUST-TRADES: 5

💰 GESAMT ROI: +31.690%
📈 Ø GEWINN PRO TRADE: +652%
📉 Ø VERLUST PRO TRADE: -182%

🏆 BESTE PERFORMANCE: +5.261% (EPIC)"""

    # === SLIDE 4: TOP PERFORMER ===
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "🏆 TOP 5 PERFORMER DER WOCHE"
    content.text = """1. 🥇 EPIC: +5.261% (Long)
2. 🥈 EPIC: +2.858% (Long) 
3. 🥉 FLOKI: +1.453% (Long)
4. 🏅 HBAR: +1.403% (Long)
5. 🏅 ENA: +1.258% (Long)

Weitere Highlights:
• PEPE: +978%
• LINK: +953% 
• GALA: +870%
• AVAX: +845%
• LDO: +824%"""

    # === SLIDE 5: COOK BEISPIEL ===
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Titel
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "📈 COOK/USDT: +142.92% ROI"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Signal Info
    signal_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(2))
    signal_frame = signal_box.text_frame
    signal_frame.text = """📱 TELEGRAM SIGNAL

Coin: COOK/USDT
Direction: Long
Leverage: 12.5X
Entry: 0.008438
Targets: 0.008522 - unlimited

✅ Signal automatisch erkannt
✅ Trade sofort ausgeführt"""
    signal_frame.paragraphs[0].font.size = Pt(14)

    # Ergebnis Info
    result_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.2), Inches(4), Inches(2))
    result_frame = result_box.text_frame
    result_frame.text = """🎯 BYBIT ERGEBNIS

Entry Price: 0.008438
Current Price: 0.009415
ROI: +142.92%

🚀 Target erreicht!
💰 Automatischer Gewinn
⏱️ Keine manuelle Arbeit"""
    result_frame.paragraphs[0].font.size = Pt(14)

    # Pfeil
    arrow = slide.shapes.add_connector(1, Inches(4.5), Inches(2.2), Inches(5.5), Inches(2.2))

    # === SLIDE 6: LTC BEISPIEL ===
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "📈 LTC/USDT: +159.00% ROI"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    signal_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(2))
    signal_frame = signal_box.text_frame
    signal_frame.text = """📱 TELEGRAM SIGNAL

Coin: LTC/USDT
Direction: Long
Leverage: 50.0X
Entry: 116.61
Targets: 117.78 - unlimited

✅ Signal automatisch erkannt
✅ Trade sofort ausgeführt"""
    signal_frame.paragraphs[0].font.size = Pt(14)

    result_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.2), Inches(4), Inches(2))
    result_frame = result_box.text_frame
    result_frame.text = """🎯 BYBIT ERGEBNIS

Entry Price: 116.61
Current Price: 120.50
ROI: +159.00%

🚀 Target übertroffen!
💰 Maximaler Hebel genutzt
⚡ Blitzschnelle Ausführung"""
    result_frame.paragraphs[0].font.size = Pt(14)

    arrow = slide.shapes.add_connector(1, Inches(4.5), Inches(2.2), Inches(5.5), Inches(2.2))

    # === SLIDE 7: TAO BEISPIEL ===
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "📈 TAO/USDT: +132.91% ROI"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    signal_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(2))
    signal_frame = signal_box.text_frame
    signal_frame.text = """📱 TELEGRAM SIGNAL

Coin: TAO/USDT
Direction: Long
Leverage: 50.0X
Entry: 434.64
Targets: 438.99 - unlimited

✅ Signal automatisch erkannt
✅ Trade sofort ausgeführt"""
    signal_frame.paragraphs[0].font.size = Pt(14)

    result_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.2), Inches(4), Inches(2))
    result_frame = result_box.text_frame
    result_frame.text = """🎯 BYBIT ERGEBNIS

Entry Price: 434.64
Current Price: 446.76
ROI: +132.91%

🚀 Perfekte Ausführung!
💰 Hoher Hebel = Hoher Gewinn
🎯 Ziel erreicht"""
    result_frame.paragraphs[0].font.size = Pt(14)

    arrow = slide.shapes.add_connector(1, Inches(4.5), Inches(2.2), Inches(5.5), Inches(2.2))

    # === SLIDE 8: ZORA MONSTER GEWINN ===
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "🔥 ZORA/USDT: +1538.61% ROI!"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)  # Rot für Aufmerksamkeit

    signal_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(2.5))
    signal_frame = signal_box.text_frame
    signal_frame.text = """📱 TELEGRAM SIGNAL

Coin: ZORA/USDT
Direction: Long
Leverage: 25.0X
Entry: 0.022240

🎯 MONSTER SIGNAL!
✅ Bot erkannte Potenzial
✅ Maximaler Hebel aktiviert
⚡ Sofortige Ausführung"""
    signal_frame.paragraphs[0].font.size = Pt(14)

    result_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.2), Inches(4), Inches(2.5))
    result_frame = result_box.text_frame
    result_frame.text = """🚀 BYBIT ERGEBNIS

Entry Price: 0.022240
Current Price: 0.036256
ROI: +1538.61%

💥 JACKPOT TRADE!
💰 Über 1500% Gewinn
🏆 Beste Performance
🤖 Bot-Perfektion"""
    result_frame.paragraphs[0].font.size = Pt(14)
    result_frame.paragraphs[0].font.color.rgb = RGBColor(0, 128, 0)  # Grün für Gewinn

    # Dicker Pfeil für Monster-Trade
    arrow = slide.shapes.add_connector(1, Inches(4.5), Inches(2.2), Inches(5.5), Inches(2.2))

    # === SLIDE 9: WARUM DIESER BOT ===
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "💡 WARUM DIESER BOT?"
    content.text = """🎯 BEWIESENE PERFORMANCE
• 84.7% Winrate in nur einer Woche
• Über 31.000% ROI
• Konsistente Gewinne

⚡ VOLLAUTOMATISCH
• Keine Emotionen beim Trading
• 24/7 Marktüberwachung
• Sofortige Signalausführung

🛡️ RISIKOMANAGEMENT
• Automatische Stop-Loss Einstellungen
• Optimierte Positionsgrößen
• Diversifizierte Trades

💰 PASSIVES EINKOMMEN
• Kein Zeitaufwand
• Keine Trading-Kenntnisse nötig
• Läuft im Hintergrund"""

    # === SLIDE 10: CALL TO ACTION ===
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "🚀 BEREIT ZUM STARTEN?"
    content.text = """📞 KONTAKT & SETUP

✅ Sofortiger Zugang zum Bot
✅ Komplette Einrichtung inklusive
✅ 24/7 Support
✅ Bewiesene Ergebnisse

💬 Telegram: @YourUsername
📧 Email: your.email@example.com
🌐 Website: www.yourbot.com

⚠️ Limitierte Plätze verfügbar!
Sichere dir jetzt deinen Zugang zum
profitabelsten Trading-Bot 2025!

🎁 Referral Code: 9EEDV9
💰 Über $5,000 in Boni verfügbar!"""

    # Präsentation speichern
    prs.save('Trading_Bot_Final.pptx')
    print("✅ FERTIG! Präsentation gespeichert als: Trading_Bot_Final.pptx")
    print("\n🎉 DEINE PRÄSENTATION ENTHÄLT:")
    print("• 10 professionelle Slides")
    print("• Echte Trading-Beispiele mit ROI")
    print("• Überzeugende Statistiken")
    print("• Call-to-Action mit Referral Code")
    print("\n📋 NÄCHSTE SCHRITTE:")
    print("1. Öffne Trading_Bot_Final.pptx")
    print("2. Passe Kontaktdaten in letzter Folie an")
    print("3. Optional: Füge weitere Beispiele hinzu")
    print("4. Präsentiere und überzeuge! 🚀")

if __name__ == "__main__":
    create_trading_presentation()
